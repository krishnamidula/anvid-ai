import base64
from datetime import datetime
from io import BytesIO
from typing import Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


class PPTXBuilder:
    COLORS = {
        "background": "#0F1117",
        "card": "#1E2130",
        "card_alt": "#171A27",
        "border": "#2D3250",
        "primary": "#6C63FF",
        "secondary": "#43B89C",
        "accent": "#FF6584",
        "warm": "#FFB347",
        "muted": "#A0A0B0",
        "text": "#FFFFFF",
        "danger": "#FF6584",
    }

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    @staticmethod
    def _rgb(hex_color: str) -> RGBColor:
        value = hex_color.replace("#", "")
        return RGBColor(int(value[:2], 16), int(value[2:4], 16), int(value[4:], 16))

    def set_dark_background(self, slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self._rgb(self.COLORS["background"])

    def add_accent_bar(self, slide, color=None, height=0.04):
        shape = slide.shapes.add_shape(1, 0, 0, self.prs.slide_width, Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(color or self.COLORS["primary"])
        shape.line.fill.background()

    def add_slide_number(self, slide, number: int, total: int):
        self.add_text(slide, f"{number}/{total}", 12.15, 7.03, 0.8, 0.25, 8, False, self.COLORS["muted"], PP_ALIGN.RIGHT)

    def add_text(self, slide, text, left, top, width, height, size=12, bold=False, color=None, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = str(text or "")
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = self._rgb(color or self.COLORS["text"])
        return box

    def add_card(self, slide, left, top, width, height, fill=None, line=None):
        shape = slide.shapes.add_shape(5, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(fill or self.COLORS["card"])
        shape.line.color.rgb = self._rgb(line or self.COLORS["border"])
        return shape

    def add_image_from_base64(self, slide, b64, left, top, width, height):
        if not b64:
            return None
        return slide.shapes.add_picture(BytesIO(base64.b64decode(b64)), Inches(left), Inches(top), Inches(width), Inches(height))

    @staticmethod
    def fmt(number):
        value = float(number or 0)
        if value >= 1_000_000:
            return f"{value / 1_000_000:.1f}M"
        if value >= 1_000:
            return f"{value / 1_000:.1f}K"
        return str(int(value))

    @staticmethod
    def _short(text, limit=150):
        value = str(text or "")
        return value if len(value) <= limit else value[:limit - 3] + "..."

    def title(self, slide, text, subtitle=None):
        self.add_text(slide, text, 0.6, 0.36, 11.5, 0.42, 23, True)
        if subtitle:
            self.add_text(slide, subtitle, 0.62, 0.82, 11.4, 0.26, 9.5, False, self.COLORS["muted"])

    def add_table(self, slide, rows, left, top, widths, row_h=0.36, winner=None):
        for r_index, row in enumerate(rows):
            x = left
            fill = self.COLORS["primary"] if r_index == 0 else self.COLORS["card"] if r_index % 2 else self.COLORS["card_alt"]
            if winner and r_index > 0 and row[0] == winner:
                self.add_card(slide, left - 0.08, top + r_index * row_h, 0.04, row_h, self.COLORS["secondary"], self.COLORS["secondary"])
            for c_index, cell in enumerate(row):
                self.add_card(slide, x, top + r_index * row_h, widths[c_index], row_h, fill, self.COLORS["border"])
                color = self.COLORS["text"] if r_index == 0 else self.COLORS["muted"]
                self.add_text(slide, cell, x + 0.05, top + r_index * row_h + 0.08, widths[c_index] - 0.1, row_h - 0.08, 7.6 if r_index else 8.2, r_index == 0, color)
                x += widths[c_index]

    def slide_cover(self, companies):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.set_dark_background(slide)
        for index in range(18):
            band = slide.shapes.add_shape(1, 0, Inches(index * 0.42), self.prs.slide_width, Inches(0.43))
            shade = 16 + index
            band.fill.solid()
            band.fill.fore_color.rgb = RGBColor(shade, shade + 2, shade + 8)
            band.line.fill.background()
        circle = slide.shapes.add_shape(9, Inches(9.2), Inches(1.05), Inches(3.7), Inches(3.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = self._rgb(self.COLORS["primary"])
        circle.fill.transparency = 30
        circle.line.fill.background()
        self.add_text(slide, "VIDEO INTELLIGENCE REPORT", 0.78, 2.1, 7.8, 0.7, 40, True)
        self.add_text(slide, "  |  ".join([company["name"] for company in companies]), 0.84, 3.02, 10.5, 0.45, 18, True, self.COLORS["secondary"])
        self.add_text(slide, f"Powered by VidStrat AI  |  {datetime.now().strftime('%d %b %Y')}", 0.84, 6.83, 6.4, 0.25, 11, False, self.COLORS["muted"])

    def build_report(self, payload: Dict) -> str:
        data = payload["companies"]
        charts = payload["charts"]
        self.slide_cover(data)
        self.slide_verdict(payload)
        self.slide_channel_overview(payload, charts)
        self.slide_content_performance(payload)
        self.slide_pillars(payload, charts)
        self.slide_cadence(payload, charts)
        self.slide_engagement(payload, charts)
        self.slide_competitor_gaps(payload)
        self.slide_recommendations(payload)
        self.slide_scorecard(payload, charts)
        for index, slide in enumerate(self.prs.slides, start=1):
            self.add_slide_number(slide, index, len(self.prs.slides))
        output = BytesIO()
        self.prs.save(output)
        output.seek(0)
        return base64.b64encode(output.read()).decode("utf-8")

    def slide_verdict(self, payload):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.set_dark_background(slide)
        self.add_accent_bar(slide)
        self.title(slide, "EXECUTIVE SUMMARY")
        self.add_text(slide, payload["executive_verdict"], 0.76, 1.28, 7.2, 4.5, 14, False)
        winner = next(company for company in payload["companies"] if company["name"] == payload["winner"])
        self.add_card(slide, 8.55, 1.2, 3.95, 4.35, self.COLORS["card"], self.COLORS["warm"])
        self.add_text(slide, "INTELLIGENCE WINNER", 8.88, 1.5, 2.8, 0.28, 12, True, self.COLORS["warm"])
        self.add_text(slide, winner["name"], 8.88, 1.9, 3.2, 0.56, 25, True)
        self.add_text(slide, f"{winner['radar_scores']['overall_strength']:.1f}", 8.88, 2.62, 1.5, 0.65, 34, True, self.COLORS["secondary"])
        reasons = [f"Engagement: {winner['radar_scores']['engagement']}", f"Performance: {winner['radar_scores']['content_performance']}", f"Activity: {winner['radar_scores']['activity_level']}"]
        self.add_text(slide, "\n".join(reasons), 8.92, 3.62, 3.0, 1.1, 11, False, self.COLORS["muted"])
        line = slide.shapes.add_shape(1, Inches(0.75), Inches(6.18), Inches(11.8), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = self._rgb(self.COLORS["secondary"])
        line.line.fill.background()

    def slide_scorecard(self, payload, charts):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.set_dark_background(slide)
        self.add_accent_bar(slide)
        self.title(slide, "SUMMARY & RANKING")
        self.add_image_from_base64(slide, charts["radar_base64"], 1.75, 1.0, 9.8, 4.15)
        rows = [["Company", "Consistency", "Engagement", "Performance", "Activity", "Overall"]]
        for item in payload["final_ranking"]:
            company = next(c for c in payload["companies"] if c["name"] == item["company"])
            score = company["radar_scores"]
            rows.append([company["name"], score["consistency"], score["engagement"], score["content_performance"], score["activity_level"], score["overall_strength"]])
        self.add_table(slide, rows, 0.8, 5.55, [2.7, 1.2, 1.2, 1.2, 1.2, 1.2], 0.32)

    def slide_channel_overview(self, payload, charts):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.set_dark_background(slide)
        self.add_accent_bar(slide)
        self.title(slide, "CHANNEL OVERVIEW")
        rows = [["Company", "Subscribers", "Total Videos", "Avg Views", "Upload Interval", "Consistency", "Eng. Rate"]]
        for company in payload["companies"]:
            rows.append([
                company["name"],
                self.fmt(company["subscribers"]),
                company["video_count"],
                self.fmt(company["avg_views"]),
                f"{company.get('avg_interval_days', 0):.1f} days",
                company["consistency_classification"],
                f"{company['avg_engagement_rate']:.2f}%"
            ])
        self.add_table(slide, rows, 0.4, 1.05, [2.3, 1.25, 1.0, 1.15, 1.2, 1.2, 1.1], 0.38, payload["winner"])
        self.add_image_from_base64(slide, charts["scatter_base64"], 2.0, 4.1, 8.8, 2.4)

    def slide_pillars(self, payload, charts):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.set_dark_background(slide)
        self.add_accent_bar(slide)
        self.title(slide, "TOPICS & THEMES")
        self.add_image_from_base64(slide, charts["pillars_base64"], 0.55, 1.15, 7.9, 4.35)
        self.add_text(slide, "WHAT THEY COVER", 9.05, 1.15, 2.6, 0.3, 13, True, self.COLORS["secondary"])
        y = 1.65
        for company in payload["companies"]:
            top = sorted(company["content_pillars"].items(), key=lambda item: item[1], reverse=True)[:2]
            self.add_text(slide, company["name"], 9.05, y, 2.9, 0.25, 10, True)
            self.add_text(slide, ", ".join([item[0] for item in top]) or "Limited pillar signal", 9.05, y + 0.24, 3.1, 0.36, 9, False, self.COLORS["muted"])
            y += 0.78
        self.add_text(slide, "MISSING THEMES", 9.05, y + 0.1, 2.6, 0.3, 13, True, self.COLORS["secondary"])
        whitespace = payload.get("whitespace_opportunities", [])[:3]
        for idx, item in enumerate(whitespace):
            self.add_text(slide, f"• {self._short(item['title'], 40)}", 9.05, y + 0.4 + idx * 0.28, 2.9, 0.3, 9, False, self.COLORS["muted"])

    def slide_content_performance(self, payload):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.set_dark_background(slide)
        self.add_accent_bar(slide)
        self.title(slide, "CONTENT PERFORMANCE")
        positions = [(0.55, 1.15), (4.55, 1.15), (8.55, 1.15), (2.55, 4.08), (6.55, 4.08)]
        for index, company in enumerate(payload["companies"][:5]):
            x, y = positions[index]
            top = (company["top_videos"] or [{}])[0]
            self.add_card(slide, x, y, 3.45, 2.25, self.COLORS["card"], [self.COLORS["primary"], self.COLORS["secondary"], self.COLORS["accent"], self.COLORS["warm"], "#A78BFA"][index])
            self.add_text(slide, company["name"], x + 0.18, y + 0.16, 2.9, 0.24, 10, True, self.COLORS["secondary"])
            self.add_text(slide, self._short(top.get("title", ""), 110), x + 0.18, y + 0.52, 3.05, 0.62, 10, True)
            self.add_text(slide, f"{self.fmt(top.get('views', 0))} views  |  {top.get('engagement_rate', 0):.2f}% ER", x + 0.18, y + 1.26, 2.8, 0.22, 8.5, False, self.COLORS["muted"])
            self.add_text(slide, self._short(top.get("why_it_worked", ""), 140), x + 0.18, y + 1.58, 3.0, 0.5, 8.2, False, self.COLORS["muted"])

    def slide_steal(self, payload):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.set_dark_background(slide)
        self.add_accent_bar(slide)
        self.title(slide, "STEAL THIS STRATEGY")
        y = 1.22
        for company in payload["companies"][:5]:
            self.add_text(slide, company["name"], 0.75, y, 2.3, 0.28, 12, True, self.COLORS["secondary"])
            self.add_text(slide, "\n".join(company["steal_strategy"][:3]), 3.0, y, 8.9, 0.65, 9, False, self.COLORS["muted"])
            y += 1.06

    def slide_cadence(self, payload, charts):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.set_dark_background(slide)
        self.add_accent_bar(slide)
        self.title(slide, "POSTING CADENCE")
        self.add_image_from_base64(slide, charts["cadence_base64"], 0.55, 1.2, 7.15, 4.9)
        y = 1.28
        for company in payload["companies"]:
            self.add_card(slide, 8.1, y, 3.85, 0.72, self.COLORS["card"], self.COLORS["border"])
            self.add_text(slide, company["name"], 8.32, y + 0.14, 1.7, 0.24, 10, True)
            self.add_text(slide, f"{company['consistency_score']:.1f}", 10.18, y + 0.08, 0.9, 0.36, 19, True, self.COLORS["secondary"])
            self.add_text(slide, company["consistency_classification"], 11.0, y + 0.18, 0.8, 0.25, 8, False, self.COLORS["muted"])
            y += 0.86

    def slide_seo(self, payload):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.set_dark_background(slide)
        self.add_accent_bar(slide)
        self.title(slide, "SEO INTELLIGENCE")
        rows = [["Company", "Overall", "Title", "Desc", "Best Day", "Timing"]]
        for company in payload["companies"]:
            seo = company["seo_scores"]
            rows.append([company["name"], seo["overall"], seo["title_score"], seo["description_score"], seo["best_upload_day"], seo["timing_score"]])
        self.add_table(slide, rows, 0.75, 1.2, [2.1, 0.9, 0.9, 0.9, 1.5, 0.9], 0.4)
        winner = max(payload["companies"], key=lambda company: company["seo_scores"]["overall"])
        self.add_card(slide, 8.6, 1.22, 3.3, 1.35, self.COLORS["card"], self.COLORS["secondary"])
        self.add_text(slide, "SEO WINNER", 8.9, 1.52, 2.0, 0.24, 10, True, self.COLORS["secondary"])
        self.add_text(slide, winner["name"], 8.9, 1.88, 2.6, 0.34, 17, True)
        self.add_text(slide, "Key insight: the fastest SEO lift will come from deeper descriptions, sharper intent-led titles, and publishing rhythm discipline.", 0.86, 4.55, 11.1, 1.05, 12, False, self.COLORS["muted"])

    def slide_engagement(self, payload, charts):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.set_dark_background(slide)
        self.add_accent_bar(slide)
        self.title(slide, "ENGAGEMENT ANALYSIS")
        self.add_image_from_base64(slide, charts["scatter_base64"], 0.55, 1.15, 6.85, 4.25)
        rows = [["Company", "Avg Views", "Avg Likes", "Avg Comments", "ER"]]
        for company in payload["companies"]:
            rows.append([company["name"], self.fmt(company["avg_views"]), self.fmt(company["avg_likes"]), self.fmt(company["avg_comments"]), f"{company['avg_engagement_rate']:.2f}%"])
        self.add_table(slide, rows, 7.65, 1.3, [1.45, 1.05, 1.0, 1.15, 0.7], 0.36)
        self.add_text(slide, "Interpretation: brands above the engagement midpoint earn stronger audience affinity, while high-view brands below it have reach that can be converted into comments, likes, and subscriber growth.", 0.75, 5.85, 11.6, 0.55, 11, False, self.COLORS["muted"])

    def slide_competitor_gaps(self, payload):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.set_dark_background(slide)
        self.add_accent_bar(slide, self.COLORS["secondary"])
        self.title(slide, "GAP ANALYSIS", "Where competitors are under-indexed")
        for index, item in enumerate(payload["whitespace_opportunities"][:5]):
            x = 0.75 + (index % 2) * 6.0
            y = 1.25 + (index // 2) * 1.75
            width = 5.35 if index < 4 else 11.35
            self.add_card(slide, x, y, width, 1.25, self.COLORS["card"], self.COLORS["secondary"])
            self.add_text(slide, item["title"], x + 0.22, y + 0.16, width - 0.4, 0.25, 11, True, self.COLORS["secondary"])
            self.add_text(slide, self._short(item["description"], 190), x + 0.22, y + 0.48, width - 0.45, 0.45, 8.6, False, self.COLORS["muted"])
            self.add_text(slide, f"{item['format']}  |  Potential: {item['potential']}", x + 0.22, y + 0.98, width - 0.45, 0.18, 8, True, self.COLORS["warm"])

    def slide_recommendations(self, payload):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.set_dark_background(slide)
        self.add_accent_bar(slide)
        self.title(slide, "RECOMMENDATIONS")
        for index, key in enumerate(["month_1", "month_2", "month_3"]):
            month = payload["action_plan"][key]
            x = 0.65 + index * 4.15
            color = [self.COLORS["primary"], self.COLORS["secondary"], self.COLORS["warm"]][index]
            self.add_card(slide, x, 1.14, 3.55, 3.35, self.COLORS["card"], color)
            self.add_text(slide, f"MONTH {index + 1}", x + 0.2, 1.36, 2.2, 0.24, 9, True, color)
            self.add_text(slide, month["title"], x + 0.2, 1.7, 2.8, 0.28, 13, True)
            self.add_text(slide, "\n".join(month["actions"][:4]), x + 0.24, 2.2, 3.0, 1.8, 8.8, False, self.COLORS["muted"])
        self.add_text(slide, "VidStrat AI", 10.45, 6.88, 2.0, 0.24, 9, False, self.COLORS["muted"], PP_ALIGN.RIGHT)
